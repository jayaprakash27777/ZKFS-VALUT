"""
ZKFS — Enhanced Professional Presentation Generator
IIT Jammu Final Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG         = RGBColor(10,  10,  15)      # Near-black bg
BG2        = RGBColor(18,  18,  28)      # Slightly lighter panel
VIOLET     = RGBColor(139, 92,  246)     # Primary violet
VIOLET_L   = RGBColor(167, 127, 253)     # Light violet for accents
CYAN       = RGBColor(34,  211, 238)     # Teal accent
EMERALD    = RGBColor(52,  211, 153)     # Success green
AMBER      = RGBColor(251, 191, 36)      # Warning/highlight amber
WHITE      = RGBColor(255, 255, 255)
ZINC_200   = RGBColor(228, 228, 231)
ZINC_400   = RGBColor(161, 161, 170)
ZINC_600   = RGBColor(82,  82,  91)
BORDER     = RGBColor(50,  50,  70)

W = Inches(13.333)
H = Inches(7.5)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, transparency=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = Pt(0)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline_textbox(slide, lines, left, top, width, height,
                           font_size=16, color=ZINC_200, spacing=Pt(6)):
    """lines = list of (text, bold, color) or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, clr = item, False, color
        else:
            text, bold, clr = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = clr
        run.font.name = "Calibri"
    return txBox

def add_divider(slide, top, color=VIOLET, width_fraction=0.15):
    """Thin horizontal accent line"""
    bar_w = W * width_fraction
    bar = slide.shapes.add_shape(1, Inches(0.6), top, bar_w, Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_accent_circle(slide, cx, cy, r, color=VIOLET):
    """Decorative circle"""
    shape = slide.shapes.add_shape(9, cx - r, cy - r, r*2, r*2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_lines, 
             accent=VIOLET, icon="●"):
    """A styled card panel"""
    # Card background
    card = add_rect(slide, left, top, width, height,
                    fill_color=BG2, border_color=BORDER)
    # Left accent bar
    bar = slide.shapes.add_shape(1, left, top, Pt(3), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Icon circle
    ci = Inches(0.22)
    ico_box = slide.shapes.add_shape(9, left + Inches(0.15), top + Inches(0.15), ci, ci)
    ico_box.fill.solid()
    ico_box.fill.fore_color.rgb = accent
    ico_box.line.fill.background()

    # Icon text
    add_text_box(slide, icon,
                 left + Inches(0.12), top + Inches(0.09),
                 ci + Inches(0.1), ci + Inches(0.1),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Card title
    add_text_box(slide, title,
                 left + Inches(0.5), top + Inches(0.12),
                 width - Inches(0.65), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)

    # Card body
    add_multiline_textbox(slide, body_lines,
                          left + Inches(0.5), top + Inches(0.52),
                          width - Inches(0.65), height - Inches(0.65),
                          font_size=11, color=ZINC_400, spacing=Pt(2))
    return card

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # Big dark gradient panel on left 55%
    add_rect(slide, 0, 0, W * 0.58, H, fill_color=RGBColor(14, 14, 22))

    # Violet glow blob (large circle behind)
    glow = slide.shapes.add_shape(9, Inches(-1.2), Inches(1.5), Inches(5), Inches(5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(88, 28, 220)
    glow.line.fill.background()
    # make it semi-transparent via xml
    sp_pr = glow.fill._xPr
    sp_pr.getparent().attrib.pop('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', None)

    # Right panel: dark card
    add_rect(slide, W * 0.6, Inches(1.2), W * 0.37, Inches(5.1),
             fill_color=BG2, border_color=BORDER)

    # IIT Jammu label (top left)
    add_text_box(slide, "IIT JAMMU", Inches(0.55), Inches(0.3), Inches(4), Inches(0.4),
                 font_size=11, bold=True, color=VIOLET, align=PP_ALIGN.LEFT)

    # Thin top accent bar
    bar = slide.shapes.add_shape(1, 0, 0, W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VIOLET
    bar.line.fill.background()

    # Main Title
    add_text_box(slide, "Zero-Knowledge", Inches(0.55), Inches(1.1), Inches(7), Inches(1.0),
                 font_size=56, bold=True, color=WHITE)
    add_text_box(slide, "File System", Inches(0.55), Inches(2.0), Inches(7), Inches(0.9),
                 font_size=56, bold=True, color=VIOLET)

    add_divider(slide, Inches(3.1), color=CYAN, width_fraction=0.18)

    add_text_box(slide, "Next-Generation Privacy-Preserving Cloud Storage",
                 Inches(0.55), Inches(3.25), Inches(7), Inches(0.5),
                 font_size=18, color=ZINC_400)

    add_text_box(slide,
                 "Presented by: Jayaprakash V\nStudent, IIT Jammu\n2025–26",
                 Inches(0.55), Inches(4.2), Inches(6), Inches(1.2),
                 font_size=14, color=ZINC_400)

    # Right info card items
    info = [
        ("🔐", "Zero-Knowledge", "Server sees only ciphertext"),
        ("⚡", "Real-Time Streaming", "AES-256-GCM chunk encryption"),
        ("🔑", "WebAuthn Passkeys", "Hardware-backed security"),
        ("☁️", "Cloud Native", "Spring Boot + MinIO + PostgreSQL"),
    ]
    cy = Inches(1.55)
    for icon, title, desc in info:
        add_text_box(slide, icon, W * 0.62, cy, Inches(0.45), Inches(0.38),
                     font_size=17, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, W * 0.62 + Inches(0.45), cy, Inches(4), Inches(0.25),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, desc, W * 0.62 + Inches(0.45), cy + Inches(0.25), Inches(4.2), Inches(0.2),
                     font_size=10, color=ZINC_400)
        cy += Inches(1.05)

    # Bottom bar
    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — TABLE OF CONTENTS
# ─────────────────────────────────────────────────────────────────────────────
def slide_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, Inches(3.8), H, fill_color=RGBColor(14, 14, 22))

    add_text_box(slide, "TABLE OF CONTENTS", Inches(0.5), Inches(0.35), Inches(3), Inches(0.4),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Agenda", Inches(0.5), Inches(0.7), Inches(3), Inches(0.65),
                 font_size=32, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=CYAN, width_fraction=0.1)

    sections = [
        ("01", "Problem Statement",        "Why existing cloud storage fails"),
        ("02", "What is ZKFS?",             "Zero-knowledge architecture overview"),
        ("03", "Security Deep Dive",        "Cryptographic mechanisms"),
        ("04", "Upload Pipeline",           "Real-time streaming encryption"),
        ("05", "Core Features",             "File sharing, passkeys, offline vault"),
        ("06", "Technology Stack",          "Frontend, Backend, Storage, Crypto"),
        ("07", "UI/UX Design",              "Premium design system"),
        ("08", "System Architecture",       "Component diagram"),
        ("09", "Future Scope",              "Roadmap and enhancements"),
    ]
    cy = Inches(1.65)
    for num, title, desc in sections:
        # Number bubble
        bubble = slide.shapes.add_shape(9, Inches(0.45), cy, Inches(0.35), Inches(0.35))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = VIOLET
        bubble.line.fill.background()
        add_text_box(slide, num, Inches(0.43), cy + Pt(2), Inches(0.4), Inches(0.3),
                     font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, Inches(0.9), cy, Inches(2.7), Inches(0.23),
                     font_size=12, bold=True, color=WHITE)
        add_text_box(slide, desc, Inches(0.9), cy + Inches(0.23), Inches(2.7), Inches(0.18),
                     font_size=9, color=ZINC_400)
        cy += Inches(0.56)

    # Right: decorative
    add_text_box(slide, "🛡️", Inches(5.5), Inches(1.5), Inches(4), Inches(4),
                 font_size=180, align=PP_ALIGN.CENTER, color=VIOLET)
    add_text_box(slide, "End-to-End Encrypted\nFile Storage Platform",
                 Inches(4.5), Inches(5.5), Inches(8), Inches(0.8),
                 font_size=18, bold=True, color=ZINC_200, align=PP_ALIGN.CENTER)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — PROBLEM STATEMENT
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "01  /  PROBLEM STATEMENT", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Why Traditional Cloud Storage Fails", Inches(0.6), Inches(0.6), Inches(10), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=AMBER, width_fraction=0.15)

    # 3-column problems
    problems = [
        ("⚠️", "Server Holds Keys", AMBER,
         "Providers like Google Drive and Dropbox manage encryption keys on their servers. They can access your files at any time."),
        ("🕵️", "Data Surveillance", RGBColor(239, 68, 68),
         "Files are scanned for advertising, compliance checks, and can be handed over to governments upon legal request."),
        ("💥", "Breach Catastrophe", RGBColor(249, 115, 22),
         "A single server compromise exposes millions of users' plaintext files. 2022–24 saw 4 billion+ records leaked from cloud providers."),
    ]
    col_w = Inches(3.8)
    cx = Inches(0.5)
    for icon, title, accent, body in problems:
        card_h = Inches(4.6)
        add_rect(slide, cx, Inches(1.7), col_w, card_h, fill_color=BG2, border_color=BORDER)
        topbar = slide.shapes.add_shape(1, cx, Inches(1.7), col_w, Pt(3))
        topbar.fill.solid(); topbar.fill.fore_color.rgb = accent
        topbar.line.fill.background()
        add_text_box(slide, icon, cx, Inches(2.0), col_w, Inches(0.7), font_size=36, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, cx, Inches(2.8), col_w, Inches(0.4),
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body, cx + Inches(0.25), Inches(3.3), col_w - Inches(0.5), Inches(2.5),
                     font_size=12, color=ZINC_400, align=PP_ALIGN.CENTER)
        cx += col_w + Inches(0.35)

    # Bottom stat banner
    add_rect(slide, 0, H - Inches(1.0), W, Inches(0.9), fill_color=RGBColor(14, 14, 22))
    stats = ["  4B+  Records leaked (2022–24)", "  93%  Of breaches involve cloud storage", "  $4.45M  Average cost of a data breach"]
    sx = Inches(0.5)
    for s in stats:
        add_text_box(slide, s, sx, H - Inches(0.85), Inches(4), Inches(0.5),
                     font_size=13, bold=True, color=AMBER)
        sx += Inches(4.4)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.17), W, Inches(0.15),
                 font_size=8, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — WHAT IS ZKFS
# ─────────────────────────────────────────────────────────────────────────────
def slide_what_is(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "02  /  OVERVIEW", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "What is ZKFS?", Inches(0.6), Inches(0.6), Inches(7), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=VIOLET, width_fraction=0.12)

    add_text_box(slide,
        "ZKFS is a Zero-Knowledge cloud storage platform where encryption and decryption happen exclusively "
        "inside the user's browser. The server stores only cryptographic ciphertext — never plaintext data.",
        Inches(0.6), Inches(1.6), Inches(7.5), Inches(0.85),
        font_size=14, color=ZINC_200)

    # Key principles as horizontal cards
    principles = [
        ("🔐", "Client-Side Only", VIOLET,
         "All cryptographic operations occur exclusively in the browser via WebCrypto API. No plaintext ever leaves your device."),
        ("🗝️", "Key Sovereignty",  CYAN,
         "You own your keys. Master keys are derived from your password using Argon2id. Passkeys use hardware-backed WebAuthn PRF."),
        ("👁️", "Server Blindness", EMERALD,
         "The server is mathematically unable to decrypt your files. Even with full database access, contents remain inaccessible."),
        ("📦", "Chunked Streaming", AMBER,
         "Files are sliced into 5MB chunks and encrypted in real-time. This supports massive files without memory overflow."),
    ]

    cy = Inches(2.65)
    for icon, title, accent, body in principles:
        h = Inches(1.05)
        add_rect(slide, Inches(0.5), cy, Inches(12.3), h, fill_color=BG2, border_color=BORDER)
        lbar = slide.shapes.add_shape(1, Inches(0.5), cy, Pt(4), h)
        lbar.fill.solid(); lbar.fill.fore_color.rgb = accent; lbar.line.fill.background()
        add_text_box(slide, icon, Inches(0.7), cy + Inches(0.25), Inches(0.6), Inches(0.5), font_size=22)
        add_text_box(slide, title, Inches(1.4), cy + Inches(0.1), Inches(2.5), Inches(0.35),
                     font_size=14, bold=True, color=WHITE)
        add_text_box(slide, body, Inches(1.4), cy + Inches(0.45), Inches(10.9), Inches(0.45),
                     font_size=11, color=ZINC_400)
        cy += Inches(1.15)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — SECURITY DEEP DIVE
# ─────────────────────────────────────────────────────────────────────────────
def slide_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "03  /  SECURITY", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Cryptographic Security Architecture", Inches(0.6), Inches(0.6), Inches(10), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=CYAN, width_fraction=0.2)

    # Left: key hierarchy diagram (text-based)
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.4), fill_color=BG2, border_color=BORDER)
    add_text_box(slide, "KEY HIERARCHY", Inches(0.7), Inches(1.75), Inches(5.4), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)

    steps = [
        ("Master Password  +  PBKDF2 Salt",  ZINC_400,  False),
        ("▼ Argon2id (4MB memory, 3 iter)",   ZINC_600,  False),
        ("Master Session Key (MSK)",            WHITE,     True),
        ("▼ AES-256-GCM Wrap",                 ZINC_600,  False),
        ("Data Encryption Key (DEK)",          CYAN,      True),
        ("▼ File encrypted chunk-by-chunk",    ZINC_600,  False),
        ("Encrypted Ciphertext  →  MinIO",     EMERALD,   True),
    ]
    ky = Inches(2.2)
    for text, color, bold in steps:
        add_text_box(slide, text, Inches(1.2), ky, Inches(4.8), Inches(0.35),
                     font_size=12, bold=bold, color=color)
        ky += Inches(0.52)

    # Right: spec cards
    specs = [
        ("Argon2id KDF",    VIOLET,  ["Memory: 4 MB per hash", "Iterations: 3 passes", "Parallelism: 2 threads", "→ Brute-force resistant"]),
        ("AES-256-GCM",     CYAN,    ["Key size: 256 bits", "Nonce: 96-bit random IV", "Auth tag: 128-bit GCM tag", "→ Authenticated encryption"]),
        ("WebAuthn PRF",    EMERALD, ["FIDO2 hardware authenticator", "PRF extension output as KEK", "Biometric-bound secret", "→ Phishing-proof passkeys"]),
        ("ECDH Sharing",    AMBER,   ["P-384 key pair per user", "Ephemeral DEK wrapping", "Secure file sharing", "→ Forward secrecy"]),
    ]
    cx = Inches(6.8)
    cy_s = Inches(1.6)
    for i, (title, accent, lines) in enumerate(specs):
        col = i % 2
        row = i // 2
        sx = cx + col * Inches(3.25)
        sy = cy_s + row * Inches(2.7)
        add_rect(slide, sx, sy, Inches(3.1), Inches(2.5), fill_color=BG2, border_color=BORDER)
        topbar = slide.shapes.add_shape(1, sx, sy, Inches(3.1), Pt(3))
        topbar.fill.solid(); topbar.fill.fore_color.rgb = accent; topbar.line.fill.background()
        add_text_box(slide, title, sx + Inches(0.15), sy + Inches(0.15), Inches(2.8), Inches(0.35),
                     font_size=13, bold=True, color=accent)
        iy = sy + Inches(0.55)
        for line in lines:
            add_text_box(slide, line, sx + Inches(0.2), iy, Inches(2.7), Inches(0.3),
                         font_size=11, color=ZINC_200 if "→" not in line else EMERALD,
                         bold=("→" in line))
            iy += Inches(0.38)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — UPLOAD PIPELINE
# ─────────────────────────────────────────────────────────────────────────────
def slide_upload(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "04  /  UPLOAD PIPELINE", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Real-Time Streaming Encryption", Inches(0.6), Inches(0.6), Inches(10), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=EMERALD, width_fraction=0.18)

    pipeline = [
        ("1", "Select File", VIOLET,
         "User selects file via drag-and-drop.\nFile stays in browser memory only."),
        ("2", "Generate DEK", CYAN,
         "Browser generates a 256-bit AES-GCM DEK\nper file using window.crypto.getRandomValues()."),
        ("3", "Wrap DEK", EMERALD,
         "DEK is wrapped with Master Key / Passkey / Password.\nServer receives only the ciphertext wrapper."),
        ("4", "Encrypt Filename", AMBER,
         "Filename and thumbnail are encrypted\nwith the Master Key for vault visibility."),
        ("5", "Chunk + Encrypt", VIOLET,
         "File is sliced into 5MB chunks. Each chunk is\nencrypted in-memory with AES-256-GCM."),
        ("6", "Upload to MinIO", CYAN,
         "Presigned URLs are used to upload encrypted chunks\ndirectly to object storage, bypassing the app server."),
        ("7", "Finalize + Verify", EMERALD,
         "SHA-256 hash of each chunk is verified.\nServer confirms upload integrity."),
    ]

    # Arrow flow horizontal
    box_w = Inches(1.6)
    box_h = Inches(2.4)
    arrow_w = Inches(0.3)
    start_x = Inches(0.3)
    mid_y = Inches(2.3)

    cx = start_x
    for i, (num, title, accent, body) in enumerate(pipeline):
        # Box
        add_rect(slide, cx, mid_y, box_w, box_h, fill_color=BG2, border_color=BORDER)
        # Top accent
        topbar = slide.shapes.add_shape(1, cx, mid_y, box_w, Pt(3))
        topbar.fill.solid(); topbar.fill.fore_color.rgb = accent; topbar.line.fill.background()
        # Number
        nb = slide.shapes.add_shape(9, cx + Inches(0.6), mid_y + Inches(0.15), Inches(0.4), Inches(0.4))
        nb.fill.solid(); nb.fill.fore_color.rgb = accent; nb.line.fill.background()
        add_text_box(slide, num, cx + Inches(0.58), mid_y + Inches(0.18), Inches(0.45), Inches(0.3),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_text_box(slide, title, cx + Inches(0.08), mid_y + Inches(0.7), box_w - Inches(0.15), Inches(0.4),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Body
        add_text_box(slide, body, cx + Inches(0.08), mid_y + Inches(1.15), box_w - Inches(0.12), Inches(1.1),
                     font_size=9, color=ZINC_400, align=PP_ALIGN.CENTER)

        cx += box_w
        if i < len(pipeline) - 1:
            # Arrow
            add_text_box(slide, "›", cx, mid_y + Inches(0.9), arrow_w, Inches(0.5),
                         font_size=22, bold=True, color=ZINC_600, align=PP_ALIGN.CENTER)
            cx += arrow_w

    # Bottom note
    add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
             fill_color=RGBColor(14, 14, 22), border_color=BORDER)
    add_text_box(slide,
        "⚡  Zero plaintext leaves the browser.  The server only ever receives AES-256-GCM encrypted bytes. "
        "Even MinIO storage operators cannot read your files.",
        Inches(0.7), Inches(5.08), Inches(12.0), Inches(0.45),
        font_size=11, color=EMERALD, bold=True)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — CORE FEATURES
# ─────────────────────────────────────────────────────────────────────────────
def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "05  /  CORE FEATURES", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Platform Capabilities", Inches(0.6), Inches(0.6), Inches(9), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=VIOLET, width_fraction=0.12)

    features = [
        ("🔐", "Passkey Authentication", VIOLET,
         ["WebAuthn FIDO2 hardware keys (Touch ID, Face ID, YubiKey)",
          "PRF extension extracts a cryptographic secret from the authenticator",
          "Used to wrap/unwrap individual file DEKs", "Phishing-proof — credential is bound to this domain"]),
        ("🔗", "Encrypted File Sharing", CYAN,
         ["P-384 ECDH ephemeral key exchange per share link",
          "DEK is re-wrapped with the recipient's public key",
          "Server sees only re-encrypted blobs", "Time-limited share tokens with expiry"]),
        ("📱", "Offline Vault Mode", EMERALD,
         ["Full vault can be exported as an encrypted .zkfs archive",
          "Entire file index and chunks are bundled locally",
          "Decrypted offline using local session keys",
          "Re-sync on reconnection with delta updates"]),
        ("🛡️", "Individual File Locks", AMBER,
         ["Lock any single file with a unique custom password",
          "Or lock with a Passkey for hardware protection",
          "DEK re-encrypted with file-specific KEK",
          "Contents require secondary auth to decrypt"]),
        ("📂", "Folder Hierarchy", VIOLET,
         ["Nested folder creation and drag-drop organization",
          "Folder names are also encrypted at rest",
          "Multi-select operations (bulk move, delete)",
          "Trash bin with 30-day soft delete"]),
        ("👥", "User Sharing", CYAN,
         ["Share files with other registered ZKFS users",
          "Asymmetric key-wrapping preserves zero-knowledge",
          "Recipient gets re-encrypted DEK, never plaintext",
          "Revokable at any time by the file owner"]),
    ]

    col_w = Inches(4.0)
    col_h = Inches(2.35)
    cols = 3
    pad_x = Inches(0.35)
    pad_y = Inches(0.25)
    sx = Inches(0.5)
    sy = Inches(1.65)

    for idx, (icon, title, accent, lines) in enumerate(features):
        row = idx // cols
        col = idx % cols
        fx = sx + col * (col_w + pad_x)
        fy = sy + row * (col_h + pad_y)

        add_rect(slide, fx, fy, col_w, col_h, fill_color=BG2, border_color=BORDER)
        lbar = slide.shapes.add_shape(1, fx, fy, Pt(3), col_h)
        lbar.fill.solid(); lbar.fill.fore_color.rgb = accent; lbar.line.fill.background()

        add_text_box(slide, icon, fx + Inches(0.1), fy + Inches(0.12),
                     Inches(0.45), Inches(0.45), font_size=20)
        add_text_box(slide, title, fx + Inches(0.6), fy + Inches(0.12),
                     col_w - Inches(0.75), Inches(0.38), font_size=12, bold=True, color=accent)

        iy = fy + Inches(0.6)
        for line in lines:
            add_text_box(slide, "• " + line, fx + Inches(0.12), iy,
                         col_w - Inches(0.2), Inches(0.28), font_size=9, color=ZINC_400)
            iy += Inches(0.38)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — TECHNOLOGY STACK
# ─────────────────────────────────────────────────────────────────────────────
def slide_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "06  /  TECHNOLOGY STACK", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "What Powers ZKFS", Inches(0.6), Inches(0.6), Inches(9), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=VIOLET, width_fraction=0.13)

    layers = [
        ("FRONTEND", VIOLET,
         [("React 18 + Next.js 14", "App Router, SSR, API routes"),
          ("TypeScript",            "End-to-end type safety"),
          ("Tailwind CSS",          "Utility-first dark theme"),
          ("Framer Motion",         "Premium micro-animations"),
          ("TanStack Query",        "Server state management")]),
        ("BACKEND", CYAN,
         [("Java 17 + Spring Boot 3", "REST API, business logic"),
          ("Spring Security + JWT",   "Stateless auth, token refresh"),
          ("Flyway",                  "SQL migration versioning"),
          ("MinIO SDK",               "Presigned S3 object upload"),
          ("Redis",                   "Session cache & revocation")]),
        ("STORAGE & DB", EMERALD,
         [("PostgreSQL 16",  "Relational metadata store"),
          ("MinIO",          "S3-compatible object storage"),
          ("Redis 7",        "Fast key-value token store"),
          ("Flyway",         "Schema version management"),
          ("JPA / Hibernate","ORM with auditing")]),
        ("CRYPTOGRAPHY", AMBER,
         [("WebCrypto API",      "Browser-native AES-256-GCM"),
          ("Argon2id (WASM)",    "Memory-hard KDF for passwords"),
          ("SimpleWebAuthn",     "FIDO2 passkey registration/auth"),
          ("ECDH P-384",        "Key exchange for file sharing"),
          ("SHA-256 Integrity", "Per-chunk hash verification")]),
    ]

    col_w = Inches(3.0)
    sx = Inches(0.4)
    for layer_title, accent, items in layers:
        add_rect(slide, sx, Inches(1.65), col_w, Inches(5.3), fill_color=BG2, border_color=BORDER)
        topbar = slide.shapes.add_shape(1, sx, Inches(1.65), col_w, Inches(0.4))
        topbar.fill.solid(); topbar.fill.fore_color.rgb = accent; topbar.line.fill.background()
        add_text_box(slide, layer_title, sx + Inches(0.1), Inches(1.7), col_w - Inches(0.15), Inches(0.3),
                     font_size=11, bold=True, color=WHITE)
        iy = Inches(2.15)
        for name, desc in items:
            add_text_box(slide, name, sx + Inches(0.15), iy, col_w - Inches(0.25), Inches(0.28),
                         font_size=11, bold=True, color=accent)
            add_text_box(slide, desc, sx + Inches(0.15), iy + Inches(0.28), col_w - Inches(0.25), Inches(0.24),
                         font_size=9, color=ZINC_400)
            iy += Inches(0.72)
        sx += col_w + Inches(0.25)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — SYSTEM ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "08  /  SYSTEM ARCHITECTURE", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Component Architecture Diagram", Inches(0.6), Inches(0.6), Inches(9), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=EMERALD, width_fraction=0.18)

    # Row 1: User Browser
    add_rect(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.8),
             fill_color=RGBColor(14, 14, 40), border_color=RGBColor(80, 60, 200))
    add_text_box(slide, "USER BROWSER (CLIENT-SIDE CRYPTO ZONE)",
                 Inches(0.55), Inches(1.7), Inches(5), Inches(0.3),
                 font_size=9, bold=True, color=VIOLET)
    comps_browser = ["Next.js UI  |  Framer Motion", "WebCrypto API  |  Argon2id WASM", "SimpleWebAuthn  |  TanStack Query", "Session Key Store (memory-only)"]
    bx = Inches(0.6)
    for comp in comps_browser:
        add_rect(slide, bx, Inches(2.05), Inches(2.8), Inches(1.1), fill_color=BG2, border_color=BORDER)
        add_text_box(slide, comp, bx + Inches(0.1), Inches(2.2), Inches(2.6), Inches(0.6),
                     font_size=10, color=ZINC_200, align=PP_ALIGN.CENTER)
        bx += Inches(3.1)

    # Arrows down
    for ax in [Inches(1.8), Inches(4.9), Inches(8.0), Inches(11.1)]:
        add_text_box(slide, "↓", ax, Inches(3.55), Inches(0.5), Inches(0.4),
                     font_size=14, color=ZINC_600, align=PP_ALIGN.CENTER)

    # Row 2: Backend
    add_rect(slide, Inches(0.4), Inches(3.95), Inches(12.5), Inches(1.3),
             fill_color=RGBColor(12, 28, 24), border_color=RGBColor(30, 130, 100))
    add_text_box(slide, "SPRING BOOT 3 BACKEND  (Java 17)",
                 Inches(0.55), Inches(4.0), Inches(5), Inches(0.3),
                 font_size=9, bold=True, color=EMERALD)
    be_comps = ["Spring Security + JWT", "FileController / AuthController", "FileStorageService + MinIO SDK", "JPA / PostgreSQL ORM"]
    bx = Inches(0.6)
    for comp in be_comps:
        add_text_box(slide, comp, bx, Inches(4.35), Inches(2.8), Inches(0.55),
                     font_size=10, color=ZINC_200, align=PP_ALIGN.CENTER)
        if bx < Inches(9):
            add_text_box(slide, "│", bx + Inches(2.8), Inches(4.35), Inches(0.3), Inches(0.55),
                         font_size=14, color=ZINC_600, align=PP_ALIGN.CENTER)
        bx += Inches(3.1)

    # Row 3: Storage
    for ax in [Inches(1.8), Inches(4.9), Inches(8.0)]:
        add_text_box(slide, "↓", ax, Inches(5.35), Inches(0.5), Inches(0.4),
                     font_size=14, color=ZINC_600, align=PP_ALIGN.CENTER)

    store_items = [("PostgreSQL 16", VIOLET, "Metadata, users,\nshare links"), ("MinIO Object Store", AMBER, "Encrypted file\nchunks (S3)"), ("Redis 7", CYAN, "JWT token cache\n& revocation")]
    sx = Inches(0.5)
    for name, accent, desc in store_items:
        add_rect(slide, sx, Inches(5.75), Inches(3.9), Inches(1.3), fill_color=BG2, border_color=BORDER)
        topb = slide.shapes.add_shape(1, sx, Inches(5.75), Inches(3.9), Pt(3))
        topb.fill.solid(); topb.fill.fore_color.rgb = accent; topb.line.fill.background()
        add_text_box(slide, name, sx + Inches(0.15), Inches(5.85), Inches(3.6), Inches(0.35),
                     font_size=12, bold=True, color=accent)
        add_text_box(slide, desc, sx + Inches(0.15), Inches(6.25), Inches(3.6), Inches(0.55),
                     font_size=10, color=ZINC_400)
        sx += Inches(4.3)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — FUTURE SCOPE
# ─────────────────────────────────────────────────────────────────────────────
def slide_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_text_box(slide, "09  /  FUTURE SCOPE", Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
                 font_size=10, bold=True, color=VIOLET)
    add_text_box(slide, "Roadmap & Enhancements", Inches(0.6), Inches(0.6), Inches(9), Inches(0.75),
                 font_size=34, bold=True, color=WHITE)
    add_divider(slide, Inches(1.45), color=AMBER, width_fraction=0.15)

    roadmap = [
        ("Q3 2026", "📱 Mobile App", VIOLET,
         "Native iOS & Android apps with biometric passkey unlock and offline encrypted vault sync."),
        ("Q4 2026", "📝 Collaborative Editing", CYAN,
         "Real-time encrypted document co-editing using CRDTs (Conflict-free Replicated Data Types). Server sees only deltas — never plaintext."),
        ("Q1 2027", "⛓️ Blockchain Audit Log", EMERALD,
         "Immutable, tamper-proof audit trail of file access events stored on a permissioned blockchain layer."),
        ("Q2 2027", "🌍 Federated Vaults", AMBER,
         "Peer-to-peer vault federation. Users can self-host ZKFS nodes and inter-operate securely with the main network."),
        ("Q3 2027", "🤖 AI on Encrypted Data", VIOLET,
         "Privacy-preserving ML inference using Fully Homomorphic Encryption (FHE). AI processes ciphertext — never plaintext."),
        ("Q4 2027", "🏢 Enterprise SSO", CYAN,
         "SAML2/OIDC federation for enterprise SSO with zero-knowledge guarantees. Group policy key management for organizations."),
    ]

    col_w = Inches(3.85)
    cx = Inches(0.5)
    cy = Inches(1.65)
    for i, (quarter, title, accent, body) in enumerate(roadmap):
        row = i // 3
        col = i % 3
        fx = cx + col * (col_w + Inches(0.3))
        fy = cy + row * Inches(2.55)
        add_rect(slide, fx, fy, col_w, Inches(2.35), fill_color=BG2, border_color=BORDER)
        lbar = slide.shapes.add_shape(1, fx, fy, Pt(3), Inches(2.35))
        lbar.fill.solid(); lbar.fill.fore_color.rgb = accent; lbar.line.fill.background()
        add_text_box(slide, quarter, fx + Inches(0.15), fy + Inches(0.1), col_w - Inches(0.3), Inches(0.25),
                     font_size=9, bold=True, color=accent)
        add_text_box(slide, title, fx + Inches(0.15), fy + Inches(0.37), col_w - Inches(0.3), Inches(0.38),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, body, fx + Inches(0.15), fy + Inches(0.82), col_w - Inches(0.3), Inches(1.4),
                     font_size=10, color=ZINC_400)

    add_text_box(slide, "ZKFS  ·  Final Project Presentation  ·  2026",
                 0, H - Inches(0.4), W, Inches(0.35),
                 font_size=9, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
def slide_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Full-width violet top stripe
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = VIOLET; bar.line.fill.background()

    # Center glow
    glow = slide.shapes.add_shape(9, Inches(2.5), Inches(0.8), Inches(8.3), Inches(8.3))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(60, 20, 130)
    glow.line.fill.background()

    add_text_box(slide, "🛡️", Inches(5.1), Inches(0.8), Inches(3), Inches(2.5),
                 font_size=100, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Thank You", Inches(1), Inches(3.1), Inches(11.3), Inches(1.3),
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Questions & Answers",
                 Inches(1), Inches(4.3), Inches(11.3), Inches(0.55),
                 font_size=22, color=ZINC_400, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(5.05), color=VIOLET, width_fraction=0.5)
    # center the divider
    for shape in slide.shapes:
        if shape.left == Inches(0.6) and shape.width < Inches(2):
            shape.left = Inches(4.0)

    add_text_box(slide, "Presented by: Jayaprakash V   ·   IIT Jammu   ·   2025–26",
                 Inches(1), Inches(5.3), Inches(11.3), Inches(0.45),
                 font_size=14, color=ZINC_400, align=PP_ALIGN.CENTER)

    add_text_box(slide, "github.com/jayaprakash27777/ZKFS-VAULT",
                 Inches(1), Inches(5.8), Inches(11.3), Inches(0.4),
                 font_size=12, color=VIOLET, align=PP_ALIGN.CENTER)

    # Bottom bar
    add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), fill_color=RGBColor(14, 14, 22))
    add_text_box(slide, "ZKFS — Zero-Knowledge File System  |  IIT Jammu  |  2026",
                 0, H - Inches(0.47), W, Inches(0.38),
                 font_size=10, color=ZINC_600, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_toc(prs)
    slide_problem(prs)
    slide_what_is(prs)
    slide_security(prs)
    slide_upload(prs)
    slide_features(prs)
    slide_tech_stack(prs)
    slide_architecture(prs)
    slide_future(prs)
    slide_thankyou(prs)

    out = "ZKFS_Enhanced_Presentation.pptx"
    prs.save(out)
    print(f"Saved -> {out} ({prs.slides.__len__()} slides)")

if __name__ == "__main__":
    main()
